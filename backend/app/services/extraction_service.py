"""Text extraction service for PDF, DOCX, and PPTX files."""

import logging
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from io import BytesIO

logger = logging.getLogger(__name__)


class ExtractionService:
    """Extracts text from various document types."""

    def extract(self, file_data: bytes, file_type: str) -> tuple[str, int]:
        """Extract text and page count from a document.
        
        Returns:
            Tuple of (extracted_text, page_count)
        """
        if file_type == "pdf":
            return self._extract_pdf(file_data)
        elif file_type == "docx":
            return self._extract_docx(file_data)
        elif file_type == "pptx":
            return self._extract_pptx(file_data)
        elif file_type in ("txt", "md"):
            return self._extract_text(file_data)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _extract_pdf(self, file_data: bytes) -> tuple[str, int]:
        """Extract text from PDF."""
        try:
            reader = PdfReader(BytesIO(file_data))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            full_text = "\n\n".join(pages)
            logger.info(f"Extracted {len(reader.pages)} pages from PDF")
            return full_text, len(reader.pages)
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise

    def _extract_docx(self, file_data: bytes) -> tuple[str, int]:
        """Extract text from DOCX."""
        try:
            doc = DocxDocument(BytesIO(file_data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)
            logger.info(f"Extracted {len(paragraphs)} paragraphs from DOCX")
            return full_text, 1
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise

    def _extract_pptx(self, file_data: bytes) -> tuple[str, int]:
        """Extract text from PPTX."""
        try:
            prs = Presentation(BytesIO(file_data))
            slides = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                slides.append("\n".join(slide_text))
            full_text = "\n\n".join(slides)
            logger.info(f"Extracted {len(prs.slides)} slides from PPTX")
            return full_text, len(prs.slides)
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise

    def _extract_text(self, file_data: bytes) -> tuple[str, int]:
        """Extract text from TXT or MD."""
        text = file_data.decode("utf-8", errors="ignore")
        return text, 1


def get_extraction_service() -> ExtractionService:
    """Return an ExtractionService instance."""
    return ExtractionService()
